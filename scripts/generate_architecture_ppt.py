from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(247, 249, 252)
TITLE = RGBColor(18, 31, 53)
TEXT = RGBColor(52, 73, 94)
MUTED = RGBColor(107, 123, 140)
LOCAL = RGBColor(227, 242, 253)
VPS = RGBColor(232, 245, 233)
EXTERNAL = RGBColor(255, 243, 224)
MCP = RGBColor(243, 229, 245)
ALERT = RGBColor(255, 235, 238)
LINE = RGBColor(110, 125, 143)


def add_box(slide, left, top, width, height, title, items, fill, title_size=20, body_size=12):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.2)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = TITLE
    p.alignment = PP_ALIGN.LEFT

    for item in items:
        p = text_frame.add_paragraph()
        p.text = f"- {item}"
        p.font.size = Pt(body_size)
        p.font.color.rgb = TEXT
        p.level = 0

    return shape


def add_label(slide, left, top, width, height, text, size=12, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return box


def connect(slide, x1, y1, x2, y2, label=None, label_left=None, label_top=None):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = LINE
    line.line.width = Pt(1.8)
    if label:
        add_label(slide, label_left, label_top, Inches(1.7), Inches(0.35), label, size=10, color=MUTED)
    return line


def build_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG

    add_label(
        slide,
        Inches(0.35),
        Inches(0.18),
        Inches(6.6),
        Inches(0.45),
        "업비트 자동매매 시스템 아키텍처",
        size=26,
        bold=True,
        color=TITLE,
    )
    add_label(
        slide,
        Inches(7.5),
        Inches(0.22),
        Inches(5.2),
        Inches(0.3),
        "로컬 연구 + VPS 실행 + MCP 계층 + 모니터링",
        size=12,
        color=MUTED,
    )

    add_box(
        slide,
        Inches(0.35),
        Inches(0.95),
        Inches(2.75),
        Inches(2.2),
        "사용자 / 로컬 노트북",
        [
            "PM Orchestrator와만 소통",
            "전략 연구 및 백테스트",
            "코드 개발 / Git / 배포",
            "대시보드 및 로그 확인",
        ],
        LOCAL,
    )

    add_box(
        slide,
        Inches(3.55),
        Inches(0.95),
        Inches(5.1),
        Inches(2.65),
        "VPS 실행 영역",
        [
            "종목 스캐너 및 신호 엔진",
            "주문 엔진 / 포지션 관리자",
            "리스크 가드 / 손실 제한",
            "WebSocket 리스너 / 스케줄러",
            "API Key / .env / 런타임 설정",
        ],
        VPS,
    )

    add_box(
        slide,
        Inches(9.05),
        Inches(0.95),
        Inches(3.9),
        Inches(2.65),
        "외부 시스템",
        [
            "Upbit Quotation API",
            "Upbit Exchange API",
            "Upbit WebSocket",
            "Telegram / Slack 알림",
        ],
        EXTERNAL,
    )

    add_box(
        slide,
        Inches(0.7),
        Inches(4.2),
        Inches(12.0),
        Inches(2.35),
        "MCP 계층",
        [
            "market-data-mcp: 시세 / 캔들 / 호가 / 거래대금 조회",
            "strategy-backtest logic: 종목 선별 / 점수화 / 규칙 평가",
            "execution-mcp: 주문 실행 / 체결 조회 / 잔고 조회",
            "risk-guard-mcp: 1회 손실 1만원 / 1일 손실 3만원 / 중복 주문 차단",
            "monitoring-mcp: 이벤트 로그 / 장애 감지 / 알림 발송",
        ],
        MCP,
        title_size=18,
        body_size=13,
    )

    add_box(
        slide,
        Inches(9.85),
        Inches(3.85),
        Inches(2.65),
        Inches(1.0),
        "모니터링",
        [
            "체결 / 장애 / 하트비트",
            "Telegram 즉시 알림",
        ],
        ALERT,
        title_size=18,
        body_size=11,
    )

    connect(
        slide,
        Inches(3.1),
        Inches(1.95),
        Inches(3.55),
        Inches(1.95),
        "배포/제어",
        Inches(3.08),
        Inches(1.58),
    )
    connect(
        slide,
        Inches(8.65),
        Inches(1.95),
        Inches(9.05),
        Inches(1.95),
        "REST/WebSocket",
        Inches(8.45),
        Inches(1.58),
    )
    connect(
        slide,
        Inches(6.1),
        Inches(3.6),
        Inches(6.1),
        Inches(4.2),
        "도구 호출",
        Inches(5.4),
        Inches(3.72),
    )
    connect(
        slide,
        Inches(10.95),
        Inches(3.6),
        Inches(11.15),
        Inches(3.85),
        "이벤트",
        Inches(10.7),
        Inches(3.6),
    )

    footer = slide.shapes.add_textbox(Inches(0.45), Inches(6.92), Inches(12.4), Inches(0.28))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = (
        "운영 원칙: 로컬은 연구·관제, VPS는 실주문·리스크 통제. "
        "사용자는 PM Orchestrator만 상대하고, 실거래 키는 VPS 고정 IPv4 환경에서만 사용."
    )
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.LEFT

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    build_presentation(Path("workspace/reports/bitcoin-trading-system-architecture.pptx"))
